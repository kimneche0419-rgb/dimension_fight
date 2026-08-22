# -*- coding: utf-8 -*-
"""Dimension Fight 발표용 PPT 생성 스크립트 (네온 다크 테마)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 네온 팔레트 ──
BG        = RGBColor(0x0A, 0x0E, 0x1A)   # 진한 우주 배경
PANEL     = RGBColor(0x14, 0x1B, 0x30)   # 카드 패널
PANEL_LN  = RGBColor(0x23, 0x2C, 0x4A)   # 패널 테두리
CYAN      = RGBColor(0x00, 0xE5, 0xFF)
MAGENTA   = RGBColor(0xFF, 0x2B, 0xD6)
YELLOW    = RGBColor(0xFF, 0xE1, 0x4D)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BODY      = RGBColor(0xDD, 0xE3, 0xF5)
DIM       = RGBColor(0x8B, 0x93, 0xB0)
FONT      = "맑은 고딕"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

ICON = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "dimension_fight_cpp", "assets", "icon_source.png")


def add_slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def glow_bar(s, x, y, w, h, color):
    """네온 느낌의 얇은 강조 바 (겹친 사각형)"""
    outer = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Emu(9525*3), y - Emu(9525*3),
                               w + Emu(9525*6), h + Emu(9525*6))
    outer.fill.solid(); outer.fill.fore_color.rgb = BG
    outer.line.color.rgb = color; outer.line.width = Pt(0.75)
    outer.shadow.inherit = False
    inner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    inner.fill.solid(); inner.fill.fore_color.rgb = color
    inner.line.fill.background(); inner.shadow.inherit = False


def text(s, x, y, w, h, runs, size=14, color=BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    """runs: 문자열 또는 [(텍스트, 색, 볼드, 크기), ...] 리스트"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, color, bold, size)]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for line_runs in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        for t, c, b, sz in line_runs:
            r = p.add_run(); r.text = t
            r.font.name = FONT; r.font.size = Pt(sz)
            r.font.color.rgb = c; r.font.bold = b
    return tb


def slide_header(s, no, title_ko, title_en):
    glow_bar(s, Inches(0.7), Inches(0.52), Inches(0.09), Inches(0.62), CYAN)
    text(s, Inches(0.95), Inches(0.38), Inches(9.5), Inches(0.5),
         [(title_ko + "  ", WHITE, True, 30), (title_en, CYAN, False, 13)])
    text(s, Inches(0.98), Inches(0.95), Inches(9.5), Inches(0.35),
         [("0" + str(no) + " / 08", DIM, False, 11)])
    # 상단 우측 브랜드
    text(s, Inches(9.6), Inches(0.45), Inches(3.1), Inches(0.4),
         [("DIMENSION ", MAGENTA, True, 12), ("FIGHT", CYAN, True, 12)],
         align=PP_ALIGN.RIGHT)


def panel(s, x, y, w, h, accent=CYAN):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = 0.06
    p.fill.solid(); p.fill.fore_color.rgb = PANEL
    p.line.color.rgb = PANEL_LN; p.line.width = Pt(1)
    p.shadow.inherit = False
    glow_bar(s, x + Inches(0.22), y + Inches(0.24), Inches(0.5), Inches(0.045), accent)
    return p


# ════════════════ 1. 표지 ════════════════
s = add_slide()
# 배경 별 장식
import random
random.seed(42)
for _ in range(60):
    x = random.uniform(0, 13.3); y = random.uniform(0, 7.5)
    sz = random.uniform(0.02, 0.06)
    c = random.choice([CYAN, MAGENTA, YELLOW, WHITE, WHITE])
    st = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(sz), Inches(sz))
    st.fill.solid(); st.fill.fore_color.rgb = c
    st.line.fill.background(); st.shadow.inherit = False

if os.path.isfile(ICON):
    pic = s.shapes.add_picture(ICON, Inches(5.79), Inches(0.72), height=Inches(2.3))

text(s, Inches(1.5), Inches(3.15), Inches(10.33), Inches(1.0),
     [[("D", MAGENTA, True, 54), ("IMENSION ", CYAN, True, 54)],
      [("F", CYAN, True, 54), ("IGHT", MAGENTA, True, 54)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.5), Inches(5.0), Inches(10.33), Inches(0.5),
     [("— 차원을 넘어, 별을 쫓아 —", YELLOW, False, 18)],
     align=PP_ALIGN.CENTER)
text(s, Inches(1.5), Inches(5.72), Inches(10.33), Inches(0.4),
     [("C++ · SDL2 기반 2D 차원 액션 슈팅  |  온라인 멀티플레이 지원", DIM, False, 14)],
     align=PP_ALIGN.CENTER)
glow_bar(s, Inches(5.42), Inches(6.35), Inches(2.5), Inches(0.05), CYAN)
text(s, Inches(1.5), Inches(6.78), Inches(10.33), Inches(0.4),
     [("Project Presentation", DIM, False, 12)], align=PP_ALIGN.CENTER)

# ════════════════ 2. 프로젝트 개요 ════════════════
s = add_slide()
slide_header(s, 2, "프로젝트 개요", "OVERVIEW")

panel(s, Inches(0.7), Inches(1.5), Inches(11.93), Inches(1.35), MAGENTA)
text(s, Inches(1.05), Inches(1.86), Inches(11.3), Inches(0.9),
     [[("「디멘션 파이트」", WHITE, True, 17)],
      [("물리 차원과 공허 차원을 자유롭게 넘나들며 싸우는 2D 액션 슈팅 게임.", BODY, False, 14)],
      [("1인 개발 · C++/SDL2 클라이언트 + 클라우드 멀티플레이 서버를 직접 구축한 풀스택 게임 프로젝트", DIM, False, 12.5)]],
     spacing=1.25)

feats = [
    ("🔀", "차원 전환 전투", "SHIFT 키로 물리·공허 차원을 실시간 전환. 차원마다 다른 적과 약점", CYAN),
    ("🧙", "직업 × 애니 열매", "전사·차원술사 등 10종 직업, 원피스·드래곤볼 등 애니 열매 스킬 세트", MAGENTA),
    ("🌐", "온라인 멀티플레이", "매치메이킹 경유 1vs1 대전 / 협동 모드. 릴레이 서버로 실시간 연결", YELLOW),
    ("☁️", "클라우드 동기화", "회원가입·로그인, 서버 저장 세이브. 어느 PC에서든 진행 상황 이어하기", CYAN),
    ("🎰", "성장 · 수집", "가챠 뽑기, 장비 제작, 우주선 업그레이드, 레벨업 빌드 선택", MAGENTA),
]
fy = Inches(3.05)
fw = Inches(2.31)
for i, (emo, t, d, ac) in enumerate(feats):
    x = Inches(0.7) + Emu(int(fw) * i) + Inches(0.08 * i)
    panel(s, x, fy, fw, Inches(3.6), ac)
    text(s, x + Inches(0.18), fy + Inches(0.38), fw - Inches(0.36), Inches(0.6),
         [(emo, WHITE, False, 26)])
    text(s, x + Inches(0.18), fy + Inches(1.0), fw - Inches(0.36), Inches(0.65),
         [(t, ac, True, 14.5)])
    text(s, x + Inches(0.18), fy + Inches(1.55), fw - Inches(0.36), Inches(1.9),
         [(d, BODY, False, 10.5)], spacing=1.25)

# ════════════════ 3. 스토리 ════════════════
s = add_slide()
slide_header(s, 3, "스토리", "STORY")

panel(s, Inches(0.7), Inches(1.5), Inches(11.93), Inches(1.85), MAGENTA)
text(s, Inches(1.05), Inches(1.78), Inches(11.3), Inches(1.4),
     [[("차원의 경계가 무너진 우주", WHITE, True, 16)],
      [("물리 차원과 공허 차원의 충돌로 우주는 '역설(Paradox)'에 빠졌다.", BODY, False, 13)],
      [("두 차원이 뒤엉키며 태어난 괴물들이 은하를 집어삼키고, 살아남은 자들은 마지막 기지 '제로 섹터'로 피난했다.", BODY, False, 13)],
      [("당신은 차원을 넘을 수 있는 마지막 파일럿. 다섯 개의 위험한 섹터를 돌파하고 특이점의 진실을 마주하라.", YELLOW, False, 13)]],
     spacing=1.3)

chapters = [
    ("Ch.1", "Sector Zero", "무중력 전투 훈련", "차원 파일럿의 첫걸음, 기본 전투 감각 익히기", CYAN),
    ("Ch.2", "Neon Ruins", "네온 폐허 시가전", "붕괴한 도시에서 보병으로 환승해 적 진압", MAGENTA),
    ("Ch.3", "Vantablack Deep", "심해 비행", "빛이 없는 심해, 어비스 괴물의 영역 돌파", CYAN),
    ("Ch.4", "Event Horizon", "사건의 지평선", "공허 균열 속, Void God의 군대와 조우", MAGENTA),
    ("Ch.5", "Singularity", "특이점", "Abyss Lord가 지키는 최종 관문 — 우주의 운명이 걸린 결전", YELLOW),
]
fy = Inches(3.75)
fw = Inches(2.31)
for i, (ch, en, ko, d, ac) in enumerate(chapters):
    x = Inches(0.7) + Emu(int(fw) * i) + Inches(0.08 * i)
    panel(s, x, fy, fw, Inches(3.0), ac)
    text(s, x + Inches(0.18), fy + Inches(0.3), fw - Inches(0.36), Inches(0.4),
         [(ch, ac, True, 13)])
    text(s, x + Inches(0.18), fy + Inches(0.68), fw - Inches(0.36), Inches(0.75),
         [(en, WHITE, True, 13.5)])
    text(s, x + Inches(0.18), fy + Inches(1.32), fw - Inches(0.36), Inches(0.4),
         [(ko, ac, True, 11.5)])
    text(s, x + Inches(0.18), fy + Inches(1.72), fw - Inches(0.36), Inches(1.2),
         [(d, DIM, False, 10)], spacing=1.15)
    if i < 4:
        text(s, x + fw - Inches(0.05), fy + Inches(1.3), Inches(0.25), Inches(0.4),
             [("›", DIM, True, 18)])

# ════════════════ 4. 기술 스택 ════════════════
s = add_slide()
slide_header(s, 4, "기술 스택", "TECH STACK")

cols = [
    ("게임 클라이언트", CYAN, [
        ("C++17", "게임 엔진 전체 (헤더 기반 모듈 구조)"),
        ("SDL2", "렌더링 · 입력 · 사운드 (Image/Mixer/TTF)"),
        ("Winsock2", "TCP 소켓 통신 (서버 연동)"),
        ("CMake + Ninja", "빌드 시스템 (MinGW GCC 15)"),
        ("PyInstaller", "네온 로딩 화면 런처 exe 패키징"),
    ]),
    ("게임 서버", MAGENTA, [
        ("Python asyncio", "비동기 TCP 서버 3종 (인증·매칭·릴레이)"),
        ("Railway", "클라우드 서버 배포 (3서비스 분리)"),
        ("PostgreSQL", "계정·세이브 저장 (SQLite 폴백 내장)"),
        ("Render", "WebSocket 게이트웨이 (경로 라우팅)"),
        ("Local Proxy", "TCP↔WSS 브리지 (WSL/콘솔 없음)"),
    ]),
    ("보안 · 인증", YELLOW, [
        ("bcrypt", "비밀번호 해싱 (레거시 평문 자동 마이그레이션)"),
        ("HMAC-SHA256", "세션 토큰 서명 (12시간 만료)"),
        ("토큰 검증", "매칭·릴레이 서버 공유 시크릿 검증"),
        ("비동기 connect", "타임아웃·논블로킹 소켓 (UI 프리징 방지)"),
    ]),
]
cx = [Inches(0.7), Inches(4.83), Inches(8.96)]
cw = Inches(3.67)
for (title, ac, items), x in zip(cols, cx):
    panel(s, x, Inches(1.5), cw, Inches(5.4), ac)
    text(s, x + Inches(0.28), Inches(1.82), cw - Inches(0.56), Inches(0.5),
         [(title, ac, True, 17)])
    # 각 항목마다 개별 박스: 기술명(굵게) + 설명을 같은 문단에 배치해 겹침 원천 차단
    iy = 2.42
    for name, desc in items:
        text(s, x + Inches(0.28), Inches(iy), cw - Inches(0.56), Inches(0.78),
             [(name, WHITE, True, 13), ("  —  " + desc, DIM, False, 10.5)],
             spacing=1.15)
        iy += 0.87

# ════════════════ 5. 플레이 방법 ════════════════
s = add_slide()
slide_header(s, 5, "플레이 방법", "HOW TO PLAY")

panel(s, Inches(0.7), Inches(1.5), Inches(6.3), Inches(5.4), CYAN)
text(s, Inches(1.0), Inches(1.78), Inches(5.7), Inches(0.45),
     [("🎮 조작키", CYAN, True, 16)])

keys = [
    ("WASD / 방향키", "기체 이동"),
    ("SPACE", "대시 (무적 회피)"),
    ("SHIFT", "차원 전환 — 물리 ↔ 공허"),
    ("Q / E", "무기 교체"),
    ("1 ~ 6", "장착 스킬 발동"),
    ("S / C / G", "상점 / 제작소 / 애니 뽑기"),
    ("I / J / V", "스킬 관리 / 직업 선택 / 기체 색상"),
    ("ESC", "로비 복귀 (자동 저장)"),
]
tbl = s.shapes.add_table(len(keys), 2, Inches(1.0), Inches(2.35),
                         Inches(5.75), Inches(4.35)).table
tbl.columns[0].width = Inches(2.15)
tbl.columns[1].width = Inches(3.6)
for i, (k, v) in enumerate(keys):
    c0, c1 = tbl.rows[i].cells
    for c, (t, col, bold) in ((c0, (k, CYAN, True)), (c1, (v, BODY, False))):
        c.fill.solid(); c.fill.fore_color.rgb = PANEL
        c.margin_top = c.margin_bottom = Emu(27432)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = t
        r.font.name = FONT; r.font.size = Pt(12.5)
        r.font.color.rgb = col; r.font.bold = bold

panel(s, Inches(7.25), Inches(1.5), Inches(5.38), Inches(5.4), MAGENTA)
text(s, Inches(7.55), Inches(1.78), Inches(4.8), Inches(0.45),
     [("🎲 게임 진행 루프", MAGENTA, True, 16)])
loop = [
    ("1.", "챕터 선택", "5개 챕터 순차 해금 (Sector Zero → Singularity)"),
    ("2.", "전투", "적 격파 → XP 오브·골드·재화 수집, 보스전 돌입"),
    ("3.", "레벨업 빌드", "연사 / 공격력 / 체력 / 실드 / 속도 / 멀티샷 중 선택"),
    ("4.", "성장 · 수집", "상점·제작·가챠로 열매·장비·기체 강화"),
    ("5.", "온라인 대결", "로비에서 매치메이킹 → 1vs1 / 협동 모드 입장"),
]
ly = Inches(2.35)
for no, t, d in loop:
    text(s, Inches(7.55), ly, Inches(0.5), Inches(0.4), [(no, YELLOW, True, 15)])
    text(s, Inches(8.05), ly - Inches(0.02), Inches(4.35), Inches(0.4),
         [(t, WHITE, True, 13.5)])
    text(s, Inches(8.05), ly + Inches(0.30), Inches(4.35), Inches(0.55),
         [(d, DIM, False, 10)], spacing=1.1)
    ly += Inches(0.88)

# ════════════════ 6. 만들면서 배운점 ════════════════
s = add_slide()
slide_header(s, 6, "만들면서 배운점", "LESSONS LEARNED")

lessons = [
    ("🔒", "동시성 버그의 무서움", CYAN,
     "DB 커넥션을 try/finally 없이 닫으면, 쿼리 하나가 실패했을 때 잠금이 풀리지 않아 이후 모든 저장이 통째로 막혀버렸다.",
     "→ 재현 테스트로 먼저 버그를 잡아낸 뒤 고치는 습관의 중요성을 체감"),
    ("⏱️", "블로킹 소켓의 함정", MAGENTA,
     "connect()를 그냥 부르면 서버가 죽어있을 때 OS 기본값(약 21초)만큼 화면이 그대로 멈춰버렸다.",
     "→ 논블로킹 소켓 + select()로 타임아웃을 직접 제어해야 진짜 반응하는 프로그램이 된다"),
    ("🔑", "비밀번호 대신 토큰", YELLOW,
     "자동 로그인을 위해 비밀번호를 파일에 저장하면 그 파일을 훔치면 끝이었다.",
     "→ 서명된 만료 토큰(HMAC)으로 대체 — 훔쳐도 시간이 지나면 무용지물"),
    ("🧩", "실행 파일도 '트리'로 산다", CYAN,
     "PyInstaller onefile로 만든 exe는 겉보기엔 하나지만 속으로 부트스트래퍼+실제 프로세스가 따로 떠 있었다.",
     "→ 단순 종료로는 프로세스가 안 죽는다는 걸 실제로 창관리자를 보며 확인 후 트리 종료로 해결"),
    ("🈶", "한글 경로가 빌드를 막을 수도", MAGENTA,
     "C++ 컴파일러(cc1plus)가 유니코드 경로에서 알 수 없는 오류로 조용히 죽는 걸 겪었다.",
     "→ 문제를 우회하기보다 원인을 끝까지 추적하는 게 결국 더 빠른 길이었다"),
    ("🌐", "게임과 웹은 언어가 다르다", YELLOW,
     "게임 클라이언트는 TCP 소켓만 쓰는데, 무료 서버는 WebSocket만 지원했다.",
     "→ 로컬 프록시로 TCP↔WSS를 직접 번역 — '중간에서 통역'하는 설계의 유용함을 배움"),
]
fy = Inches(1.55)
fw = Inches(3.84)
fh = Inches(2.75)
for i, (emo, title, ac, prob, learn) in enumerate(lessons):
    col = i % 3; row = i // 3
    x = Inches(0.7) + Emu(int(fw) * col) + Inches(0.08 * col)
    y = fy + Emu(int(fh) * row) + Inches(0.12 * row)
    panel(s, x, y, fw, fh, ac)
    text(s, x + Inches(0.2), y + Inches(0.22), fw - Inches(0.4), Inches(0.45),
         [(emo + "  ", WHITE, False, 15), (title, ac, True, 13)])
    text(s, x + Inches(0.2), y + Inches(0.78), fw - Inches(0.4), Inches(1.15),
         [(prob, BODY, False, 10.5)], spacing=1.22)
    text(s, x + Inches(0.2), y + Inches(2.02), fw - Inches(0.4), Inches(0.68),
         [(learn, YELLOW if ac != YELLOW else CYAN, True, 10.5)], spacing=1.2)

# ════════════════ 7. 좋은점과 개선점 ════════════════
s = add_slide()
slide_header(s, 7, "좋은점과 개선점", "STRENGTHS & IMPROVEMENTS")

panel(s, Inches(0.7), Inches(1.5), Inches(5.85), Inches(5.4), CYAN)
text(s, Inches(1.0), Inches(1.78), Inches(5.25), Inches(0.45),
     [("✅ 좋은 점", CYAN, True, 16)])
goods = [
    ("완성도 높은 콘텐츠 볼륨", "5개 챕터, 10종 직업, 애니 IP 열매·스킬, 가챠·제작·강화까지 혼자서 방대한 콘텐츠 구현"),
    ("독창적인 핵심 메커닉", "차원 전환(SHIFT) 하나로 이동·전투·전략이 동시에 바뀌는 게임만의 정체성 확보"),
    ("실서비스 수준의 인프라", "회원가입·로그인·클라우드 저장·매치메이킹·릴레이까지 직접 설계한 백엔드 3종"),
    ("보안까지 챙긴 설계", "bcrypt 해싱 + HMAC 서명 토큰으로 평문 비밀번호를 어디에도 남기지 않음"),
]
gy = Inches(2.35)
for t, d in goods:
    text(s, Inches(1.0), gy, Inches(5.25), Inches(0.35), [(t, WHITE, True, 12.5)])
    text(s, Inches(1.0), gy + Inches(0.33), Inches(5.25), Inches(0.65),
         [(d, DIM, False, 10)], spacing=1.15)
    gy += Inches(1.13)

panel(s, Inches(6.78), Inches(1.5), Inches(5.85), Inches(5.4), MAGENTA)
text(s, Inches(7.08), Inches(1.78), Inches(5.25), Inches(0.45),
     [("🔧 개선하고 싶은 점", MAGENTA, True, 16)])
improves = [
    ("네트워크 지연 대응", "P2P 릴레이 구조라 플레이어 간 핑 차이에 취약함 — 서버 측 입력 동기화 구조로 보완 예정"),
    ("자동화된 테스트", "수동 재현 테스트로 버그를 잡았던 경험을 바탕으로 유닛·통합 테스트를 코드에 미리 심어두기"),
    ("밸런스 데이터 기반 튜닝", "직업·무기·챕터 난이도가 감으로 설계됨 — 플레이 로그를 모아 수치 기반으로 다듬을 계획"),
    ("온보딩 UX", "처음 접하는 사람에겐 조작·시스템이 한 번에 몰려옴 — 튜토리얼 단계 도입이 필요"),
]
iy = Inches(2.35)
for t, d in improves:
    text(s, Inches(7.08), iy, Inches(5.25), Inches(0.35), [(t, WHITE, True, 12.5)])
    text(s, Inches(7.08), iy + Inches(0.33), Inches(5.25), Inches(0.65),
         [(d, DIM, False, 10)], spacing=1.15)
    iy += Inches(1.13)

# ════════════════ 8. 인사말 ════════════════
s = add_slide()
random.seed(7)
for _ in range(50):
    x = random.uniform(0, 13.3); y = random.uniform(0, 7.5)
    sz = random.uniform(0.02, 0.06)
    c = random.choice([CYAN, MAGENTA, YELLOW, WHITE, WHITE])
    st = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(sz), Inches(sz))
    st.fill.solid(); st.fill.fore_color.rgb = c
    st.line.fill.background(); st.shadow.inherit = False

text(s, Inches(1.5), Inches(2.65), Inches(10.33), Inches(1.0),
     [[("감", MAGENTA, True, 46), ("사", CYAN, True, 46), ("합", MAGENTA, True, 46),
       ("니", CYAN, True, 46), ("다", MAGENTA, True, 46)]],
     align=PP_ALIGN.CENTER)
glow_bar(s, Inches(5.42), Inches(3.55), Inches(2.5), Inches(0.05), CYAN)
text(s, Inches(1.5), Inches(3.95), Inches(10.33), Inches(0.9),
     [("혼자서 기획부터 클라이언트, 서버, 배포까지 부딪혀가며 완성한 프로젝트입니다.", BODY, False, 14),
      ], align=PP_ALIGN.CENTER)
text(s, Inches(1.5), Inches(4.35), Inches(10.33), Inches(0.5),
     [("끝까지 들어주셔서 감사합니다. 함께 차원을 넘어봐요.", YELLOW, False, 14)],
     align=PP_ALIGN.CENTER)

panel(s, Inches(4.42), Inches(5.15), Inches(4.5), Inches(1.15), MAGENTA)
text(s, Inches(4.72), Inches(5.42), Inches(3.9), Inches(0.65),
     [("Q & A", MAGENTA, True, 15)] , align=PP_ALIGN.CENTER)
text(s, Inches(4.72), Inches(5.72), Inches(3.9), Inches(0.5),
     [("궁금하신 점, 편하게 질문해 주세요!", DIM, False, 11)], align=PP_ALIGN.CENTER)

text(s, Inches(1.5), Inches(6.7), Inches(10.33), Inches(0.4),
     [("DIMENSION FIGHT  ·  Thank You", DIM, False, 12)], align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Dimension_Fight_프레젠테이션.pptx")
prs.save(out)
print("SAVED:", out)
